import io
from typing import Any


def generate_pptx_pitch_deck_bytes(
    problem_name: str,
    dossier: dict[str, Any],
    market_analysis: dict[str, Any],
    risk_assessment: dict[str, Any],
) -> bytes:
    """
    Generates a 7-slide python-pptx pitch deck for the selected problem.
    Returns the PPTX file as raw bytes (suitable for upload to Cloudinary).

    Slides:
      1. Title
      2. Market Gap
      3. MVP Blueprint
      4. Market Size
      5. Competitive Landscape
      6. 4-Week Roadmap
      7. Risk Audit
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Brand colours
    VIOLET = RGBColor(0x6D, 0x28, 0xD9)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1E, 0x1B, 0x4B)
    GREY = RGBColor(0x6B, 0x72, 0x80)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def _add_slide() -> Any:
        return prs.slides.add_slide(blank_layout)

    def _add_text_box(
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 18,
        bold: bool = False,
        color: RGBColor = DARK,
        align: Any = PP_ALIGN.LEFT,
        word_wrap: bool = True,
    ) -> None:
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _add_bullet_text_box(
        slide: Any,
        bullets: list[str],
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 14,
        color: RGBColor = DARK,
    ) -> None:
        from pptx.util import Pt as _Pt
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for bullet in bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.level = 1
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = _Pt(font_size)
            run.font.color.rgb = color

    def _fill_slide_background(slide: Any, color: RGBColor) -> None:
        from pptx.util import Emu
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # -----------------------------------------------------------------------
    # Slide 1 — Title
    # -----------------------------------------------------------------------
    slide1 = _add_slide()
    _fill_slide_background(slide1, DARK)
    _add_text_box(
        slide1, "FluxIdeas — Founder's Dossier",
        left=0.5, top=0.5, width=12, height=1,
        font_size=16, color=RGBColor(0xC4, 0xB5, 0xFD),
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide1, problem_name,
        left=0.5, top=1.8, width=12, height=2,
        font_size=40, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide1, "Founder's Dossier by FluxIdeas",
        left=0.5, top=4.0, width=12, height=0.8,
        font_size=18, color=RGBColor(0xA7, 0x8B, 0xFA),
        align=PP_ALIGN.CENTER,
    )

    # -----------------------------------------------------------------------
    # Slide 2 — Market Gap
    # -----------------------------------------------------------------------
    market_gap_score = dossier.get("market_gap_score", {})
    gap_rationale = market_gap_score.get("rationale", "No market gap data available.")
    signal = dossier.get("signal_strength", {})
    mention_count = signal.get("mention_count", "N/A")
    source_summary = signal.get("source_summary", "")

    slide2 = _add_slide()
    _add_text_box(slide2, "Market Gap Analysis", 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color=VIOLET)
    _add_text_box(slide2, gap_rationale, 0.5, 1.3, 11.5, 2.5, font_size=16, color=DARK)
    _add_bullet_text_box(
        slide2,
        [
            f"Urgency: {market_gap_score.get('urgency', 'N/A')}/10",
            f"Commercial Potential: {market_gap_score.get('commercial_potential', 'N/A')}/10",
            f"Feasibility: {market_gap_score.get('feasibility', 'N/A')}/10",
            f"Overall Gap Score: {market_gap_score.get('total', 'N/A')}/10",
            f"Online Mentions: ~{mention_count}",
            source_summary,
        ],
        left=0.5, top=4.0, width=11.5, height=3.0,
        font_size=14,
    )

    # -----------------------------------------------------------------------
    # Slide 3 — MVP Blueprint
    # -----------------------------------------------------------------------
    mvp = dossier.get("mvp_blueprint", {})
    feature_bullets: list[str] = []
    for key in ["feature_1", "feature_2", "feature_3"]:
        feat = mvp.get(key, {})
        if isinstance(feat, dict):
            name = feat.get("name", "")
            desc = feat.get("description", "")
            if name:
                feature_bullets.append(f"{name}: {desc}")

    slide3 = _add_slide()
    _add_text_box(slide3, "MVP Blueprint", 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color=VIOLET)
    _add_bullet_text_box(slide3, feature_bullets or ["No MVP data available."], 0.5, 1.5, 11.5, 5.0, font_size=16)

    # -----------------------------------------------------------------------
    # Slide 4 — Market Size
    # -----------------------------------------------------------------------
    tam = market_analysis.get("tam", "N/A")
    sam = market_analysis.get("sam", "N/A")
    som = market_analysis.get("som", "N/A")
    growth_rate = market_analysis.get("growth_rate", "N/A")
    verdict = market_analysis.get("economist_verdict", "")

    slide4 = _add_slide()
    _add_text_box(slide4, "Market Size (TAM / SAM / SOM)", 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color=VIOLET)
    _add_bullet_text_box(
        slide4,
        [
            f"TAM — Total Addressable Market: {tam}",
            f"SAM — Serviceable Addressable Market: {sam}",
            f"SOM — Serviceable Obtainable Market: {som}",
            f"Growth Rate: {growth_rate}",
        ],
        left=0.5, top=1.5, width=11.5, height=3.5,
        font_size=18,
    )
    if verdict:
        _add_text_box(slide4, verdict, 0.5, 5.2, 11.5, 1.8, font_size=14, color=GREY)

    # -----------------------------------------------------------------------
    # Slide 5 — Competitive Landscape
    # -----------------------------------------------------------------------
    competitors = dossier.get("competitive_landscape", [])
    comp_bullets: list[str] = []
    for comp in competitors:
        if isinstance(comp, dict):
            name = comp.get("competitor", "?")
            weakness = comp.get("weakness", "")
            edge = comp.get("your_edge", "")
            comp_bullets.append(f"{name} — Weakness: {weakness} | Our Edge: {edge}")

    slide5 = _add_slide()
    _add_text_box(slide5, "Competitive Landscape", 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color=VIOLET)
    _add_bullet_text_box(slide5, comp_bullets or ["No competitor data available."], 0.5, 1.5, 11.5, 5.5, font_size=14)

    # -----------------------------------------------------------------------
    # Slide 6 — 4-Week Roadmap
    # -----------------------------------------------------------------------
    roadmap = dossier.get("technical_roadmap", {})
    week_plan = roadmap.get("week_plan", [])
    week_bullets: list[str] = []
    for week in week_plan:
        if isinstance(week, dict):
            week_bullets.append(f"Week {week.get('week', '?')}: {week.get('focus', '')}")
    tech_stack = roadmap.get("tech_stack", "")
    timeline = roadmap.get("timeline", "")

    slide6 = _add_slide()
    _add_text_box(slide6, "4-Week Technical Roadmap", 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color=VIOLET)
    if tech_stack:
        _add_text_box(slide6, f"Stack: {tech_stack}", 0.5, 1.3, 11.5, 0.5, font_size=14, color=GREY)
    if timeline:
        _add_text_box(slide6, f"Timeline: {timeline}", 0.5, 1.9, 11.5, 0.5, font_size=14, color=GREY)
    _add_bullet_text_box(slide6, week_bullets or ["No roadmap data available."], 0.5, 2.6, 11.5, 4.5, font_size=16)

    # -----------------------------------------------------------------------
    # Slide 7 — Risk Audit
    # -----------------------------------------------------------------------
    tech_risk = risk_assessment.get("technical_risk", "N/A")
    market_risk = risk_assessment.get("market_risk", "N/A")
    legal_risk = risk_assessment.get("legal_risk", "N/A")
    kill_switch = risk_assessment.get("kill_switch_criteria", "N/A")
    survival = risk_assessment.get("survival_strategy", "N/A")

    slide7 = _add_slide()
    _add_text_box(slide7, "Red-Team Risk Audit", 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color=VIOLET)
    _add_bullet_text_box(
        slide7,
        [
            f"Technical Risk: {tech_risk}",
            f"Market Risk: {market_risk}",
            f"Legal Risk: {legal_risk}",
            f"Kill Switch: {kill_switch}",
            f"Survival Strategy: {survival}",
        ],
        left=0.5, top=1.5, width=11.5, height=5.5,
        font_size=14,
    )

    # -----------------------------------------------------------------------
    # Serialise to bytes
    # -----------------------------------------------------------------------
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
